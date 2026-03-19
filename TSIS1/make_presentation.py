from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──────────────────────────────────────────────────────────
RED       = RGBColor(0xE0, 0x1C, 0x24)   # Kaspi red
DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # near-black
DARK2     = RGBColor(0x16, 0x21, 0x3E)   # dark navy
MID       = RGBColor(0x0F, 0x3C, 0x60)   # deep blue
ACCENT    = RGBColor(0x53, 0xB8, 0xD4)   # light cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xF0, 0xF4, 0xF8)
GRAY      = RGBColor(0xA0, 0xAA, 0xB8)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)
YELLOW    = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper: solid background ───────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Helper: add rectangle ─────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

# ── Helper: add textbox ───────────────────────────────────────────────────
def tb(slide, text, l, t, w, h,
       size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
       italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

# ── Helper: multi-line textbox ────────────────────────────────────────────
def tb_lines(slide, lines, l, t, w, h,
             size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             line_bold=None, line_color=None, spacing_after=None):
    """lines = list of str  OR  list of (str, bold, color)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, lb, lc = line
        else:
            txt, lb, lc = line, bold, color
        if line_bold is not None: lb = line_bold
        if line_color is not None: lc = line_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spacing_after:
            p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = lb
        run.font.color.rgb = lc

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)

# big red bar left
rect(s, 0, 0, 0.55, 7.5, RED)
# accent line
rect(s, 0.55, 2.8, 8.5, 0.06, ACCENT)

tb(s, "TSIS ASSIGNMENT", 0.8, 0.6, 10, 0.7, size=20, bold=True, color=ACCENT)
tb(s, "JTBD Interviewer Agent", 0.8, 1.2, 11, 1.4, size=44, bold=True, color=WHITE)
tb(s, "AI-Powered Legal Compliance Agent for Kaspi Bank", 0.8, 2.9, 10, 0.6, size=22, color=LIGHT_GRAY)

tb_lines(s, [
    "Cloud Computing for Big Data  ·  7 Points  ·  Week 4",
    "Company: Kaspi Bank JSC  ·  Domain: Legal Compliance",
    "Student deliverable: Strategic Portfolio (5 artifacts)"
], 0.8, 3.8, 11, 1.5, size=16, color=GRAY)

# Kaspi stats box
rect(s, 8.8, 4.6, 4.1, 2.5, MID)
rect(s, 8.8, 4.6, 4.1, 0.06, ACCENT)
tb(s, "KASPI FY 2024", 8.95, 4.65, 3.8, 0.4, size=13, bold=True, color=ACCENT)
tb_lines(s, [
    "$5.32B Revenue  (+32% YoY)",
    "$2.18B Net Income  (+25% YoY)",
    "10.2M Monthly Active Users",
    "75% Operating Margin",
], 8.95, 5.1, 3.8, 1.8, size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — ASSIGNMENT OVERVIEW
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, MID)
rect(s, 0, 1.08, 13.33, 0.05, RED)

tb(s, "ASSIGNMENT OVERVIEW", 0.4, 0.15, 10, 0.5, size=14, bold=True, color=ACCENT)
tb(s, "The Scenario & Mission", 0.4, 0.45, 10, 0.55, size=26, bold=True, color=WHITE)

# Scenario box
rect(s, 0.3, 1.3, 12.7, 1.4, MID)
tb(s, "⚡  THE SCENARIO", 0.5, 1.35, 5, 0.38, size=13, bold=True, color=ACCENT)
tb(s, "You are appointed Head of Cloud Transformation at Kaspi Bank. The Board demands 'Digital Transformation' and wants AI integrated into Legal Compliance. However the org is stuck in an IT-Service mindset, not a Product Operating Model. Your mission: diagnose the gaps, find the real problem, present a governed roadmap.",
   0.5, 1.7, 12.3, 0.9, size=14, color=WHITE)

# 3 phase boxes
phases = [
    ("PHASE 1", "Strategic\nDiagnosis", "Cagan Audit", "Audit Kaspi against\nMarty Cagan's Product\nOperating Model using\nreal Annual Reports", RED),
    ("PHASE 2", "Deep\nDiscovery", "JTBD Simulation", "Interview an AI persona\n(Kaspi Legal Officer)\nusing 5 Whys to find\nthe root cause", ACCENT),
    ("PHASE 3", "Managerial\nDefense", "Governance & Finance", "Calculate ROI, apply\nKazakhstan AI Ethics\ntest, write Go/No-Go\nexecutive summary", GREEN),
]
for i, (label, title, sub, desc, color) in enumerate(phases):
    x = 0.3 + i * 4.35
    rect(s, x, 2.9, 4.1, 4.2, MID)
    rect(s, x, 2.9, 4.1, 0.08, color)
    tb(s, label,      x+0.15, 2.95, 3.8, 0.35, size=11, bold=True, color=color)
    tb(s, title,      x+0.15, 3.28, 3.8, 0.9,  size=22, bold=True, color=WHITE)
    tb(s, sub,        x+0.15, 4.15, 3.8, 0.4,  size=13, color=color, bold=True)
    tb(s, desc,       x+0.15, 4.55, 3.8, 1.6,  size=13, color=LIGHT_GRAY)
    # phase number circle
    circ = slide_circle(s, x+3.55, 2.9, 0.52, color) if False else None

def add_badge(slide, x, y, text, bg_col):
    r = rect(slide, x, y, 0.45, 0.45, bg_col)
    tb(slide, text, x+0.03, y+0.04, 0.4, 0.38, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i,(label,*_,color) in enumerate(phases):
    add_badge(s, 0.3+i*4.35+3.55, 6.6, str(i+1), color)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PHASE 1: SYSTEM PROMPT & METHOD
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x8B, 0x0D, 0x0F))
rect(s, 0, 1.08, 13.33, 0.05, RED)

tb(s, "PHASE 1 — STRATEGIC DIAGNOSIS", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=RGBColor(0xFF,0xAA,0xAA))
tb(s, "The Cagan Audit: System Prompt & Methodology", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

# Left: what we loaded
rect(s, 0.3, 1.25, 5.8, 5.9, MID)
rect(s, 0.3, 1.25, 5.8, 0.06, RED)
tb(s, "📂  CONTEXT LOADED INTO AI", 0.5, 1.3, 5.4, 0.4, size=13, bold=True, color=RED)
tb_lines(s, [
    "1.  Kaspi FY 2024 Annual Report (SEC 20-F)",
    "     Revenue $5.32B · Margin 75% · 9,637 staff",
    "",
    "2.  4Q/FY 2024 Earnings Release (ir.kaspi.kz)",
    "     10.2M MAU · GMV +39% · Payments +40%",
    "",
    "3.  KASE Disclosure (ticker: CSBN)",
    "     3 regulators: NBK, AFDRA, FinMonitoring",
    "",
    "4.  ESG Report 2023",
    "     $5M community investment · Sustainalytics",
    "",
    "5.  Governance pages (ir.kaspi.kz)",
    "     CEO Lomtadze (Harvard MBA) · NPS-based mgmt",
], 0.5, 1.75, 5.5, 5.2, size=12.5, color=WHITE)

# Right: 4 pillars
rect(s, 6.4, 1.25, 6.6, 5.9, MID)
rect(s, 6.4, 1.25, 6.6, 0.06, RED)
tb(s, "🔍  CAGAN'S 4 PILLARS AUDITED", 6.6, 1.3, 6.2, 0.4, size=13, bold=True, color=RED)

pillars = [
    ("1", "Empowered Product Teams",   "Do back-office teams own outcomes\nor just execute tickets?",           "⚠️ PARTIAL"),
    ("2", "Outcome vs. Output",         "Does Legal measure impact or\njust contracts-reviewed count?",         "❌ FAIL"),
    ("3", "Continuous Discovery",       "Does Legal interview internal\ncustomers to find unmet needs?",         "❌ FAIL"),
    ("4", "Missionary vs. Mercenary",   "Are support teams driven by\nmission or by KPI box-checking?",         "❌ FAIL"),
]
for i, (num, title, q, status) in enumerate(pillars):
    y = 1.75 + i * 1.32
    rect(s, 6.5, y, 0.45, 0.45, RED)
    tb(s, num, 6.5, y+0.04, 0.45, 0.38, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, 7.05, y,      5.7, 0.32, size=13, bold=True, color=WHITE)
    tb(s, q,     7.05, y+0.3,  4.5, 0.55, size=11, color=LIGHT_GRAY)
    col = ORANGE if "PARTIAL" in status else RED
    tb(s, status, 11.2, y, 1.6, 0.4, size=11, bold=True, color=col, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — TRANSFORMATION MEMO (Phase 1 output)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x8B, 0x0D, 0x0F))
rect(s, 0, 1.08, 13.33, 0.05, RED)

tb(s, "PHASE 1 — OUTPUT", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=RGBColor(0xFF,0xAA,0xAA))
tb(s, "Transformation Memo: Key Findings", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

findings = [
    (RED,    "Empowered Teams — PARTIAL PASS",
             "Consumer product teams ARE empowered (NPS-driven, autonomous). Legal/Compliance are NOT — they operate as a shared service responding to tickets. No Legal Product Manager exists. No outcome ownership."),
    (RED,    "Outcome Orientation — FAIL",
             "Consumer products measure MAU and NPS. Legal measures contracts reviewed per week. With $5.32B revenue flowing through 3 regulatory bodies, this is a strategic liability."),
    (RED,    "Continuous Discovery — FAIL",
             "Zero evidence of structured internal discovery in Legal or Procurement. These teams wait to be asked — classic IT-service trap (Cagan). The legal team has never been treated as a product stakeholder."),
    (ORANGE, "Missionary Culture — FAIL (Back-Office)",
             "Consumer mission ('improve lives of Kazakhstanis through tech') is genuine. Back-office exhibits mercenary dynamics: Legal reviews to avoid liability, not to enable business. Output-optimized, not outcome-oriented."),
]
for i, (color, title, body) in enumerate(findings):
    y = 1.22 + i * 1.52
    rect(s, 0.3, y, 0.08, 1.3, color)
    rect(s, 0.45, y, 12.5, 1.3, RGBColor(0x1E, 0x2D, 0x45))
    tb(s, title, 0.65, y+0.08, 12, 0.38, size=14, bold=True, color=color)
    tb(s, body,  0.65, y+0.45, 12, 0.78, size=12.5, color=LIGHT_GRAY)

tb(s, "→  File: transformation_memo.md", 0.4, 7.1, 8, 0.3, size=11, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — PHASE 2: PERSONA & METHOD
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x06, 0x52, 0x6B))
rect(s, 0, 1.08, 13.33, 0.05, ACCENT)

tb(s, "PHASE 2 — DEEP DISCOVERY", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=ACCENT)
tb(s, "The JTBD Interview: AI Persona & 5 Whys", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

# Persona card
rect(s, 0.3, 1.22, 5.5, 3.5, MID)
rect(s, 0.3, 1.22, 5.5, 0.06, ACCENT)
tb(s, "👤  AI PERSONA", 0.5, 1.28, 5, 0.35, size=13, bold=True, color=ACCENT)
tb(s, "Aizat Bekova", 0.5, 1.65, 5, 0.5, size=22, bold=True, color=WHITE)
tb(s, "Senior Legal Counsel, Kaspi Bank", 0.5, 2.12, 5, 0.35, size=13, color=ACCENT)
tb_lines(s, [
    "7 years at Kaspi (joined 2018)",
    "Reviews 15–25 contracts per week",
    "Tools: Word + tracked changes + email",
    "Works until 9 PM regularly",
    "",
    '"My manager calls me thorough.',
    ' I would call myself terrified."',
], 0.5, 2.5, 5.1, 2.1, size=12, color=LIGHT_GRAY)

# Surface problem
rect(s, 6.1, 1.22, 6.9, 1.6, MID)
rect(s, 6.1, 1.22, 6.9, 0.06, ORANGE)
tb(s, "SURFACE PROBLEM (before 5 Whys)", 6.3, 1.28, 6.5, 0.4, size=12, bold=True, color=ORANGE)
tb(s, '"Contract review is too slow — takes 2-8 days per contract"', 6.3, 1.65, 6.5, 0.9, size=13, color=WHITE, italic=True)

# Root cause
rect(s, 6.1, 3.0, 6.9, 1.8, MID)
rect(s, 6.1, 3.0, 6.9, 0.06, GREEN)
tb(s, "ROOT CAUSE (after 5 Whys)", 6.3, 3.06, 6.5, 0.4, size=12, bold=True, color=GREEN)
tb(s, '"Lawyers are personally liable under KZ banking law — so they rationally RESIST automation that removes their control, even when they know it would help."',
   6.3, 3.44, 6.5, 1.28, size=13, color=WHITE, italic=True)

# JTBD statement
rect(s, 0.3, 4.9, 12.7, 1.35, RGBColor(0x06, 0x3A, 0x50))
rect(s, 0.3, 4.9, 12.7, 0.06, ACCENT)
tb(s, "JOB TO BE DONE", 0.5, 4.96, 5, 0.35, size=12, bold=True, color=ACCENT)
tb(s, '"When I receive a contract for review, help me quickly identify high-risk clauses so I can make confident, DEFENSIBLE decisions — without fear of personal regulatory liability."',
   0.5, 5.32, 12.3, 0.82, size=14, color=WHITE, italic=True)

tb(s, "The job is NOT 'review contracts faster.' The job is 'make confident decisions without fear.'",
   0.5, 6.35, 12, 0.4, size=13, bold=True, color=ACCENT)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — THE 5 WHYS CHAIN
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x06, 0x52, 0x6B))
rect(s, 0, 1.08, 13.33, 0.05, ACCENT)

tb(s, "PHASE 2 — INTERVIEW LOG", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=ACCENT)
tb(s, "The 5 Whys Chain — Root Cause Discovery", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

whys = [
    ("WHY 1", "Why does review take so long?",
     "No pre-screening exists. Every contract, regardless of size, reaches senior counsel with zero pre-processing. Lawyer is first AND last line of defense.",
     ACCENT),
    ("WHY 2", "Why no pre-screening tool?",
     "IT proposed 18-month implementation timeline. Budget was cut. Meanwhile Marketplace GMV grew +39% YoY, generating more contracts without more lawyers.",
     ACCENT),
    ("WHY 3", "Why hasn't Legal pushed harder for tooling?",
     "Lawyers fear: if a tool makes a mistake and they 'approved the output' — are they still protected? Under KZ law, regulatory liability names individuals, not software.",
     ORANGE),
    ("WHY 4", "Why no governance framework for AI tools?",
     "Nobody asked Legal to co-design it. AI mandate comes top-down. Legal received announcements, not frameworks. They were not in the room when tools were chosen.",
     ORANGE),
    ("WHY 5", "Why excluded from the design process?",
     "Legal is seen as overhead — a cost center that slows things down. They are treated as passive executors, not empowered experts. Yet they understand risk better than anyone.",
     RED),
]
for i, (label, q, a, color) in enumerate(whys):
    y = 1.22 + i * 1.22
    rect(s, 0.3, y, 1.1, 1.05, color)
    tb(s, label, 0.3, y+0.1, 1.1, 0.45, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, str(i+1), 0.3, y+0.48, 1.1, 0.5, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.55, y, 11.4, 1.05, RGBColor(0x1A, 0x2C, 0x44))
    tb(s, q, 1.75, y+0.05, 11, 0.38, size=13, bold=True, color=color)
    tb(s, a, 1.75, y+0.42, 11, 0.55, size=12, color=LIGHT_GRAY)

tb(s, "→  File: interview_log.md", 0.4, 7.2, 8, 0.25, size=11, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — PRD OVERVIEW
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x06, 0x52, 0x6B))
rect(s, 0, 1.08, 13.33, 0.05, ACCENT)

tb(s, "PHASE 2 — OUTPUT", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=ACCENT)
tb(s, "PRD: Kaspi Legal Intelligence Agent (KLIA)", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

# Left column: what KLIA does
rect(s, 0.3, 1.22, 6.1, 5.9, MID)
rect(s, 0.3, 1.22, 6.1, 0.06, ACCENT)
tb(s, "🤖  WHAT KLIA DOES", 0.5, 1.28, 5.7, 0.38, size=13, bold=True, color=ACCENT)
tb_lines(s, [
    "1.  Ingests contracts (PDF / Word, up to 100 pages)",
    "2.  Detects & classifies clauses by type",
    "     (liability, data processing, governing law...)",
    "3.  Cross-references KZ legal ruleset",
    "     (Civil Code · Law No. 94-V · AFDRA rules)",
    "4.  Assigns risk score: GREEN / AMBER / RED",
    "5.  Generates structured risk report with",
    "     specific legal references + suggestions",
    "6.  Logs immutable audit trail (7-year retention)",
    "7.  Routes RED contracts to Head of Legal",
    "",
    "⚠️  KLIA does NOT approve contracts.",
    "     Human lawyer has final authority — always.",
    "     This distinction is critical for liability.",
], 0.5, 1.72, 5.7, 5.2, size=12.5, color=WHITE)

# Right column: key requirements
rect(s, 6.7, 1.22, 6.3, 2.8, MID)
rect(s, 6.7, 1.22, 6.3, 0.06, ACCENT)
tb(s, "⚙️  KEY REQUIREMENTS", 6.9, 1.28, 5.9, 0.38, size=13, bold=True, color=ACCENT)
tb_lines(s, [
    "Languages:    Kazakh + Russian + English",
    "Response:     < 5 min for 30-page contract",
    "Availability: 99.5% uptime (09:00–22:00 AST)",
    "Data:         On-premises, KZ residency compliant",
    "Audit log:    Retained 7 years (financial reqs)",
    "Explainability: Every flag cites a legal article",
], 6.9, 1.7, 6.1, 2.2, size=12.5, color=WHITE)

# Success metrics
rect(s, 6.7, 4.22, 6.3, 2.9, MID)
rect(s, 6.7, 4.22, 6.3, 0.06, GREEN)
tb(s, "📈  SUCCESS METRICS", 6.9, 4.28, 5.9, 0.38, size=13, bold=True, color=GREEN)
metrics = [
    ("Review cycle time", "2-8 days", "< 1 day"),
    ("Lawyer overtime/week", "~10 hrs", "< 3 hrs"),
    ("Contracts/lawyer/week", "15-25", "40-60"),
]
for i, (metric, before, after) in enumerate(metrics):
    y = 4.72 + i * 0.72
    tb(s, metric, 6.9, y, 3.0, 0.38, size=12, color=LIGHT_GRAY)
    tb(s, before, 9.9, y, 1.2, 0.38, size=12, color=ORANGE, align=PP_ALIGN.CENTER)
    tb(s, "→", 11.1, y, 0.4, 0.38, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, after,  11.5, y, 1.3, 0.38, size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

tb(s, "→  File: prd.md", 0.4, 7.2, 8, 0.25, size=11, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — ROI CALCULATION
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x0F, 0x4F, 0x1A))
rect(s, 0, 1.08, 13.33, 0.05, GREEN)

tb(s, "PHASE 3 — FINANCIAL REASONING", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=GREEN)
tb(s, "ROI Calculation: Is KLIA Worth Building?", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

# Cost of problem
rect(s, 0.3, 1.22, 5.8, 3.1, MID)
rect(s, 0.3, 1.22, 5.8, 0.06, ORANGE)
tb(s, "💸  COST OF THE PROBLEM", 0.5, 1.28, 5.4, 0.38, size=13, bold=True, color=ORANGE)
tb_lines(s, [
    "Senior lawyer salary (hh.kz):  650,000 KZT/mo",
    "Hourly rate:  4,063 KZT/hr  (~$8.75)",
    "Wasted time:  10 hrs/week × 4 weeks = 40 hrs/mo",
    "Cost per lawyer/month:  162,500 KZT",
    "",
    "Legal team size (est.):  20 senior lawyers",
    "Total monthly waste:  3,250,000 KZT",
    "Annual cost of status quo:  39,000,000 KZT",
    "                           (~$84,000 USD/year)",
], 0.5, 1.7, 5.5, 2.5, size=12.5, color=WHITE)

# Cost of solution
rect(s, 6.4, 1.22, 6.6, 3.1, MID)
rect(s, 6.4, 1.22, 6.6, 0.06, ACCENT)
tb(s, "🔧  COST OF THE SOLUTION", 6.6, 1.28, 6.2, 0.38, size=13, bold=True, color=ACCENT)
tb_lines(s, [
    "One-time build:       26,000,000 KZT  (~$56K)",
    "  · Fine-tuning + rule engine + frontend",
    "  · On-premises API gateway (data residency)",
    "  · Legal governance sign-off",
    "",
    "Annual operating:      7,284,000 KZT  (~$16K)",
    "  · Claude API tokens:  84,000 KZT ($180/yr!)",
    "  · Maintenance 0.5 FTE: 6,000,000 KZT",
    "  · Server hosting:    1,200,000 KZT",
], 6.6, 1.7, 6.3, 2.5, size=12.5, color=WHITE)

# ROI summary
rect(s, 0.3, 4.5, 12.7, 2.65, RGBColor(0x0A, 0x38, 0x14))
rect(s, 0.3, 4.5, 12.7, 0.06, GREEN)
tb(s, "📊  ROI SUMMARY", 0.5, 4.56, 5, 0.38, size=13, bold=True, color=GREEN)

roi_items = [
    ("Net Annual Benefit (from Year 2):", "31,716,000 KZT  /  $68,300 USD", GREEN),
    ("Payback Period:", "~10 months", WHITE),
    ("Year 1 ROI:", "17.2%", YELLOW),
    ("Year 2+ ROI:", "435% annually", GREEN),
    ("3-Year NPV (12% discount):", "+32,681,000 KZT  /  ~$70,400 USD", GREEN),
    ("API token cost vs. salary savings:", "$180/year vs. $84,000/year  (0.2%)", ACCENT),
]
for i, (label, val, color) in enumerate(roi_items):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.35
    y = 5.0 + row * 0.62
    tb(s, label, x, y, 3.5, 0.38, size=12.5, color=GRAY)
    tb(s, val,   x+3.5, y, 2.7, 0.38, size=12.5, bold=True, color=color)

tb(s, "→  File: executive_summary.md", 0.4, 7.2, 8, 0.25, size=11, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — GOVERNANCE & ETHICS
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK2)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x4A, 0x1A, 0x00))
rect(s, 0, 1.08, 13.33, 0.05, ORANGE)

tb(s, "PHASE 3 — GOVERNANCE & ETHICS", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=ORANGE)
tb(s, "Kazakhstan AI Ethics Test: 5 Risks Evaluated", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

risks = [
    ("DATA RESIDENCY", "HIGH → mitigated",
     "KZ Law No. 94-V (2021): PII of KZ citizens must stay in Kazakhstan. Cannot send contracts to foreign cloud AI directly. Mitigation: on-premises API gateway.",
     ORANGE, "LOW after fix"),
    ("AI LIABILITY GAP", "HIGH → managed",
     "KZ banking law names individual officers — not software — as liable. No legal precedent for AI compliance calls. Mitigation: formal AI Decision-Support Policy + AFDRA notification.",
     RED, "MEDIUM"),
    ("HALLUCINATION RISK", "HIGH → mitigated",
     "LLMs can invent legal citations. Mitigation: deterministic citation engine (rule DB, not LLM) for all legal references. LLM used only for clause detection + explanation.",
     RED, "LOW-MED"),
    ("TRANSPARENCY", "MED → managed",
     "KZ National AI Concept 2024-2029 requires explainable AI in sensitive domains. Mitigation: every risk flag cites a specific legal article. Full audit trail exported for AFDRA.",
     YELLOW, "LOW"),
    ("VENDOR DEPENDENCY", "MED → low",
     "Single foreign AI vendor (Anthropic) is a risk for critical compliance function. Mitigation: model-agnostic API layer — swap to local/open-source model without rebuilding product.",
     YELLOW, "LOW"),
]
for i, (title, severity, body, color, residual) in enumerate(risks):
    col = i % 2 if i < 4 else 0
    row = i // 2
    w = 6.35 if i < 4 else 12.7
    x = 0.3 + col * 6.65
    y = 1.22 + row * 1.88
    rect(s, x, y, w-0.2, 1.7, RGBColor(0x1A, 0x28, 0x3D))
    rect(s, x, y, w-0.2, 0.06, color)
    tb(s, title, x+0.15, y+0.08, w-0.8, 0.35, size=12, bold=True, color=color)
    tb(s, severity, x+w-1.7, y+0.05, 1.4, 0.38, size=11, bold=True, color=color, align=PP_ALIGN.RIGHT)
    tb(s, body, x+0.15, y+0.42, w-0.5, 0.88, size=11.5, color=LIGHT_GRAY)
    tb(s, f"Residual: {residual}", x+0.15, y+1.38, w-0.5, 0.3, size=10.5, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — GO / NO-GO VERDICT
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.1, RGBColor(0x0B, 0x4A, 0x1A))
rect(s, 0, 1.08, 13.33, 0.05, GREEN)

tb(s, "PHASE 3 — EXECUTIVE DECISION", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=GREEN)
tb(s, "Go / No-Go Verdict", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

# Big verdict box
rect(s, 0.3, 1.22, 12.7, 1.5, RGBColor(0x0A, 0x3D, 0x18))
rect(s, 0.3, 1.22, 12.7, 0.08, GREEN)
rect(s, 0.3, 1.22, 2.6, 1.5, GREEN)
tb(s, "✅ GO", 0.3, 1.38, 2.6, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "CONDITIONAL", 0.3, 1.95, 2.6, 0.55, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, "Proceed with Phase 1 pilot. The financial case is strong. Governance risks are real but manageable. The worst outcome is NOT building — it means continuing to scale Legal reactively as Kaspi's Marketplace grows at 39% YoY.",
   3.1, 1.32, 9.7, 1.2, size=14, color=WHITE)

# Conditions
rect(s, 0.3, 2.95, 12.7, 3.7, MID)
rect(s, 0.3, 2.95, 12.7, 0.06, GREEN)
tb(s, "CONDITIONS FOR GO", 0.5, 3.0, 6, 0.4, size=13, bold=True, color=GREEN)
conditions = [
    ("1", "Data Residency Architecture Review — CTO to certify KLIA meets Law No. 94-V before any contract data is processed.", GREEN),
    ("2", "AI Decision-Support Policy — CLO approves and publishes internal policy defining KLIA role, human obligations, liability attribution. Due 30 days before pilot.", ACCENT),
    ("3", "AFDRA Notification — Compliance team notifies AFDRA of pilot under operational change process. Engage proactively, not reactively.", ACCENT),
    ("4", "Pilot Scope Limit — Phase 1: LOW and MEDIUM risk contracts only. RED contracts stay on manual review until 6-month pilot validated.", ORANGE),
    ("5", "Hallucination Testing — Run KLIA against 100 historical contracts with known-correct analysis. Accuracy must exceed 90% before launch.", ORANGE),
]
for i, (num, text, color) in enumerate(conditions):
    y = 3.48 + i * 0.62
    rect(s, 0.5, y+0.05, 0.38, 0.38, color)
    tb(s, num, 0.5, y+0.05, 0.38, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, text, 1.05, y+0.05, 11.8, 0.5, size=12.5, color=LIGHT_GRAY)

tb(s, "→  File: executive_summary.md", 0.4, 6.78, 8, 0.25, size=11, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — DELIVERABLES SUMMARY
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 1.1, MID)
rect(s, 0, 1.08, 13.33, 0.05, ACCENT)

tb(s, "SUBMISSION PORTFOLIO", 0.4, 0.15, 12, 0.4, size=13, bold=True, color=ACCENT)
tb(s, "All 5 Artifacts — TSIS1/ Folder", 0.4, 0.5, 12, 0.5, size=24, bold=True, color=WHITE)

artifacts = [
    ("system_prompt.txt",      "System Prompt",       "Exact instructions for Cagan Audit, loading Kaspi FY2024 Annual Report, 4Q Earnings, KASE disclosures, ESG Report. Specifies 4 pillars and output format.",                RED,    "Phase 1"),
    ("transformation_memo.md", "Transformation Memo", "1-page AI-generated audit citing $5.32B revenue, 75% margins, NPS org model. Identifies 4 failures in back-office vs. Cagan's Product Operating Model.",                         RED,    "Phase 1"),
    ("interview_log.md",       "Interview Log",        "Full transcript of 5 Whys with Aizat Bekova persona. Non-obvious insight: lawyers resist AI because of personal liability structure, not technophobia.",                         ACCENT, "Phase 2"),
    ("prd.md",                 "Product PRD",          "Full Product Requirements Document for Kaspi Legal Intelligence Agent (KLIA) — JTBD, 8 functional requirements, success metrics, technical architecture, stakeholders.",          ACCENT, "Phase 2"),
    ("executive_summary.md",   "Executive Summary",    "ROI: 435% Year 2+, $70K NPV, $180/yr API cost vs. $84K/yr savings. 5 governance risks with KZ-specific laws (No. 94-V, AFDRA). Conditional Go/No-Go verdict.", GREEN,  "Phase 3"),
    ("kaspi_sources.md",       "Source Documents",     "All 12 primary and secondary sources: SEC 20-F, KASE filings, Earnings Release, ESG Report, KZ laws cited. Full citation list with URLs and specific data points extracted.",   GRAY,   "Context"),
]
for i, (filename, title, desc, color, phase) in enumerate(artifacts):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.55
    y = 1.22 + row * 1.98
    rect(s, x, y, 6.25, 1.8, MID)
    rect(s, x, y, 6.25, 0.06, color)
    tb(s, phase,    x+5.1, y+0.08, 1.0, 0.28, size=10, color=color, align=PP_ALIGN.RIGHT)
    tb(s, title,    x+0.15, y+0.08, 5.8, 0.38, size=14, bold=True, color=WHITE)
    tb(s, filename, x+0.15, y+0.45, 5.8, 0.28, size=11, color=color, italic=True)
    tb(s, desc,     x+0.15, y+0.72, 5.9, 1.0,  size=11, color=LIGHT_GRAY)

tb(s, "Grading rubric: Context Engineering · Discovery Depth · Financial Logic · Governance & Ethics  —  5pts each", 0.3, 7.2, 13, 0.28, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════
out = "/home/bali/Рабочий стол/Cloud-computing-university/TSIS1/TSIS_JTBD_Agent_Kaspi.pptx"
prs.save(out)
print(f"Saved: {out}")
